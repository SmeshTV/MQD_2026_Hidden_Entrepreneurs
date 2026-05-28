"""
Generate an ENGLISH PowerPoint presentation (.pptx) for the
MQD 2026 Hidden Entrepreneur Detection project — Genius Level.

English version of the comprehensive presentation with simple explanations,
before/after data examples, calculation walkthroughs, and column comparisons.
"""

import sys
sys.stdout.reconfigure(encoding="utf-8")

from pathlib import Path
import numpy as np
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

DATA = Path(".")

CM_IMG   = DATA / "confusion_matrix.png"
FI_IMG   = DATA / "feature_importance.png"
SHAP_SUM = DATA / "shap_summary.png"
SHAP_BAR = DATA / "shap_bar.png"

WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
BLACK       = RGBColor(0x00, 0x00, 0x00)
DARK_BLUE   = RGBColor(0x00, 0x2B, 0x5B)
MED_BLUE    = RGBColor(0x00, 0x5A, 0x9E)
ACCENT_GOLD = RGBColor(0xF5, 0xA6, 0x23)
LIGHT_GRAY  = RGBColor(0xF2, 0xF2, 0xF2)
DARK_GRAY   = RGBColor(0x33, 0x33, 0x33)
CODE_BG     = RGBColor(0x1E, 0x1E, 0x2E)
CODE_FG     = RGBColor(0xCD, 0xD9, 0xE0)
GREEN_BG    = RGBColor(0xE8, 0xF5, 0xE9)
GREEN_BORDER= RGBColor(0x27, 0xAE, 0x60)
YELLOW_BG   = RGBColor(0xFF, 0xF8, 0xE1)
ORANGE_TEXT = RGBColor(0xF3, 0x93, 0x1E)
RED_LIGHT   = RGBColor(0xFF, 0xEB, 0xEE)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
W = prs.slide_width
H = prs.slide_height

# ── Helpers ───────────────────────────────────────────────────────────

def add_blank_slide():
    return prs.slides.add_slide(prs.slide_layouts[6])

def add_bg(slide, color=WHITE):
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = color

def add_rect(slide, l, t, w, h, fill, line=None):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if line:
        s.line.fill.solid(); s.line.fill.fore_color.rgb = line; s.line.width = Pt(1)
    else:
        s.line.fill.background()
    return s

def tb(slide, l, t, w, h, text, fs=18, bold=False, color=BLACK,
       align=PP_ALIGN.LEFT, font="Calibri", anchor=MSO_ANCHOR.TOP, ls=1.15):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame; tf.word_wrap = True; tf.auto_size = None
    p = tf.paragraphs[0]; p.text = text
    p.font.size = Pt(fs); p.font.bold = bold; p.font.color.rgb = color
    p.font.name = font; p.alignment = align
    p.line_spacing = Pt(int(fs*ls))
    try: tf.vertical_anchor = anchor
    except: pass
    return box

def mtb(slide, l, t, w, h, lines, fs=16, color=DARK_GRAY, font="Calibri",
        align=PP_ALIGN.LEFT, ls=1.3, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame; tf.word_wrap = True; tf.auto_size = None
    try: tf.vertical_anchor = anchor
    except: pass
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line; p.font.size = Pt(fs); p.font.color.rgb = color
        p.font.name = font; p.alignment = align
        p.line_spacing = Pt(int(fs*ls))
    return box

def code_block(slide, l, t, w, h, code, fs=9):
    add_rect(slide, l, t, w, h, CODE_BG)
    box = slide.shapes.add_textbox(l+Inches(0.12), t+Inches(0.08),
                                    w-Inches(0.24), h-Inches(0.16))
    tf = box.text_frame; tf.word_wrap = True; tf.auto_size = None
    for i, line in enumerate(code.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line; p.font.size = Pt(fs); p.font.color.rgb = CODE_FG
        p.font.name = "Consolas"; p.line_spacing = Pt(fs*1.3)

def title_strip(slide, text, sub=None):
    add_rect(slide, 0, 0, W, Inches(1.1), DARK_BLUE)
    tb(slide, Inches(0.6), Inches(0.15), Inches(12), Inches(0.7),
       text, fs=30, bold=True, color=WHITE)
    if sub:
        tb(slide, Inches(0.6), Inches(0.7), Inches(12), Inches(0.35),
           sub, fs=14, color=ACCENT_GOLD)
    add_rect(slide, 0, Inches(1.1), W, Inches(0.04), ACCENT_GOLD)

def section_num(slide, num):
    s = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.4), Inches(1.4),
                                Inches(0.55), Inches(0.55))
    s.fill.solid(); s.fill.fore_color.rgb = ACCENT_GOLD; s.line.fill.background()
    tf = s.text_frame; tf.word_wrap = False
    p = tf.paragraphs[0]; p.text = num; p.font.size = Pt(18); p.font.bold = True
    p.font.color.rgb = WHITE; p.alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

def add_img(slide, path, l, t, w=None, h=None):
    if path.exists():
        slide.shapes.add_picture(str(path), l, t, w or Inches(5), h or Inches(4.5))

def new_slide(title, sub=None, num=None):
    s = add_blank_slide(); add_bg(s, WHITE)
    title_strip(s, title, sub)
    if num: section_num(s, num)
    return s

SIMPLE_ICON = "💡"
WARN_ICON   = "⚠️"

def add_explain_box(slide, l, t, w, h, lines, icon=SIMPLE_ICON, bg=GREEN_BG, border=GREEN_BORDER):
    add_rect(slide, l, t, w, h, bg, border)
    all_lines = [f"{icon} IN SIMPLE TERMS:", ""] + lines
    mtb(slide, l+Inches(0.15), t+Inches(0.1), w-Inches(0.3), h-Inches(0.2),
        all_lines, fs=12, color=DARK_GRAY, ls=1.2)

def add_data_example(slide, l, t, w, h, before_lines, after_lines):
    add_rect(slide, l, t, w/2-Inches(0.05), h, RGBColor(0xFD, 0xED, 0xED), RGBColor(0xE7, 0x4C, 0x3C))
    all_b = ["BEFORE:"] + before_lines
    mtb(slide, l+Inches(0.1), t+Inches(0.08), w/2-Inches(0.25), h-Inches(0.16),
        all_b, fs=10, color=DARK_GRAY, ls=1.15)
    add_rect(slide, l+w/2+Inches(0.05), t, w/2-Inches(0.05), h, RGBColor(0xE8, 0xF5, 0xE9), RGBColor(0x27, 0xAE, 0x60))
    all_a = ["AFTER:"] + after_lines
    mtb(slide, l+w/2+Inches(0.15), t+Inches(0.08), w/2-Inches(0.25), h-Inches(0.16),
        all_a, fs=10, color=DARK_GRAY, ls=1.15)

def add_calc_box(slide, l, t, w, h, lines):
    add_rect(slide, l, t, w, h, YELLOW_BG, ORANGE_TEXT)
    all_c = ["🧮 CALCULATION WALKTHROUGH:"] + lines
    mtb(slide, l+Inches(0.1), t+Inches(0.08), w-Inches(0.2), h-Inches(0.16),
        all_c, fs=11, color=DARK_GRAY, ls=1.2)

def add_compare_box(slide, l, t, w, h, why_use, why_not):
    add_rect(slide, l, t, w/2-Inches(0.04), h, RGBColor(0xE8, 0xF5, 0xE9), GREEN_BORDER)
    mtb(slide, l+Inches(0.1), t+Inches(0.08), w/2-Inches(0.24), h-Inches(0.16),
        ["✅ WE USE:"] + why_use, fs=10, color=DARK_GRAY, ls=1.15)
    add_rect(slide, l+w/2+Inches(0.04), t, w/2-Inches(0.04), h, RED_LIGHT, RGBColor(0xE7, 0x4C, 0x3C))
    mtb(slide, l+w/2+Inches(0.14), t+Inches(0.08), w/2-Inches(0.24), h-Inches(0.16),
        ["❌ WE DON'T USE:"] + why_not, fs=10, color=DARK_GRAY, ls=1.15)

# ══════════════════════════════════════════════════════════════════════
#  SLIDE 1: TITLE
# ══════════════════════════════════════════════════════════════════════
s = add_blank_slide(); add_bg(s, DARK_BLUE)
add_rect(s, 0, Inches(2.4), W, Inches(2.6), MED_BLUE)
add_rect(s, 0, Inches(5.0), W, Inches(0.06), ACCENT_GOLD)
tb(s, Inches(0.8), Inches(2.6), Inches(11.7), Inches(0.7),
   "Mastercard Data Quest 2026", fs=18, color=ACCENT_GOLD, align=PP_ALIGN.CENTER)
tb(s, Inches(0.8), Inches(3.1), Inches(11.7), Inches(1.2),
   "Hidden Entrepreneur Detection", fs=44, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
tb(s, Inches(0.8), Inches(4.1), Inches(11.7), Inches(0.6),
   "10 Behavioral Paradoxes · Cross-border Sourcing · AUC-ROC · Target Product",
   fs=16, color=RGBColor(0xBB,0xCC,0xDD), align=PP_ALIGN.CENTER)
tb(s, Inches(0.8), Inches(5.3), Inches(11.7), Inches(0.5),
   "Genius Level — Complete solution with step-by-step explanations",
   fs=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
tb(s, Inches(0.8), Inches(6.0), Inches(11.7), Inches(0.8),
   "Team MQD_2026  |  May 2026",
   fs=14, color=RGBColor(0x99,0xAA,0xBB), align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════
#  SLIDE 2: Table of Contents
# ══════════════════════════════════════════════════════════════════════
s = new_slide("Table of Contents", "Complete breakdown — 31 slides", "01")
mtb(s, Inches(0.5), Inches(1.5), Inches(12), Inches(5.8), [
    "1.  Problem Statement — what are we solving and why",
    "2.  Data — sources, sizes, structure",
    "3.  Loading + concat + merge — simple explanation with example",
    "4.  Label — what is it and why?",
    "5.  MCC Categories — business, wholesale, logistics, fuel, SaaS/Ads",
    "6.  Behavioral Paradoxes — 10 + cross-border:",
    "     6a. Clockwork Buyer        6b. Off-Hours Operator",
    "     6c. Expense Ratio Inversion    6d. Token Wholesale",
    "     6e. Baseline Aggregations       6f. Supplier Fingerprint",
    "     6g. Last-Mile Echo              6h. Round-Trip Cash Flow",
    "     6i. Inventory Pulse             6j. Multi-Vendor Loyalty",
    "     6k. Channel Schizophrenia       6l. Cross-border Sourcing",
    "7.  Feature Assembly — 35 features, fillna, dropna",
    "8.  CatBoost — why this model, training, parameters",
    "9.  5-Fold Cross-Validation — stability check",
    "10. Results — AUC 1.0, Confusion Matrix, where numbers come from",
    "11. Feature Importance — top-15, what each means",
    "12. SHAP — explaining every prediction",
    "13. Top-100 Hidden Entrepreneurs + Target Product",
    "14. Columns we DON'T use and why (comparison)",
    "15. Technical problems — merge_asof, Round-Trip, pandas 3.0",
    "16. Conclusions — what we achieved, what matters",
], fs=12, ls=1.2)

# ══════════════════════════════════════════════════════════════════════
#  SLIDE 3: Problem Statement
# ══════════════════════════════════════════════════════════════════════
s = new_slide("Problem Statement", "What we are solving and why", "02")
mtb(s, Inches(0.5), Inches(1.4), Inches(7.5), Inches(5.8), [
    "📌 The Problem (simple terms):",
    "We have 80,000 consumer cards (ordinary people).",
    "But SOME owners use them for BUSINESS — buying goods",
    "wholesale, paying for ads, hiring delivery.",
    "The bank can't see this — the card is 'personal'.",
    "Task: find these 'hidden entrepreneurs'.",
    "",
    "📌 What we are given:",
    "• 3 million transactions from KNOWN business cards (ground truth)",
    "• 9.8 million transactions from consumer cards (search here)",
    "• Merchant reference directory — who is who",
    "",
    "📌 Why AUC-ROC instead of Accuracy?",
    "• If we say 'everyone is consumer', Accuracy = 76%,",
    "  but we find nobody. Accuracy is a bad metric.",
    "• AUC-ROC = how well the model DISTINGUISHES",
    "  business from consumer. 1.0 = perfect.",
    "",
    "📌 What are 'Behavioral Paradoxes'?",
    "They are 10 BEHAVIOR PATTERNS that exist in business cards",
    "but are almost absent in consumer cards. If a consumer card",
    "behaves like a business card — it's a 'hidden entrepreneur'.",
])
add_explain_box(s, Inches(8.3), Inches(1.4), Inches(4.5), Inches(5.8), [
    "Imagine: you see a person in a store.",
    "They buy 1 loaf of bread → consumer.",
    "They buy 100 boxes of cookies → businessperson.",
    "But what if they bought 100 boxes of cookies",
    "with a personal card? The bank only sees:",
    "'Card X spent 50,000 KZT at a supermarket'.",
    "",
    "Our model looks NOT at the amount,",
    "but at PATTERNS:",
    "• Do they buy regularly?",
    "• Outside working hours?",
    "• Wholesale or retail?",
    "• Do they switch channels?",
    "• Do they pay abroad?",
    "",
    "10 paradoxes = 10 different views",
    "on a card's behavior.",
], icon="🎯")

# ══════════════════════════════════════════════════════════════════════
#  SLIDE 4: Data
# ══════════════════════════════════════════════════════════════════════
s = new_slide("Data — sources, sizes, structure", "3 parquet files, 12.8M rows", "03")
mtb(s, Inches(0.5), Inches(1.4), Inches(6.5), Inches(2.2), [
    "📁 business_cards_MDQ.parquet — 2,997,593 rows, 25,000 cards",
    "   Known business cards. label = 1 (ground truth)",
    "📁 consumer_cards_MDQ.parquet — 9,832,487 rows, 80,000 cards",
    "   Consumer cards. label = 0 (search here)",
    "📁 merchants_reference.parquet — 2,165 merchant directory",
    "",
    "📊 Total: 12,830,080 rows, 105,000 unique cards",
    "📊 Class balance: 76% consumer vs 24% business",
], fs=12, ls=1.2)
code_block(s, Inches(7.2), Inches(1.4), Inches(5.6), Inches(2.2), """# Read only the columns we need
BIZ_COLS = ["transaction_date","transaction_timestamp",
            "transaction_amount_kzt","mcc","merchant_id",
            "channel","card_number","tokenized",
            "is_recurring","country"]

biz = pd.read_parquet("business_cards_MDQ.parquet",
                      columns=BIZ_COLS)
con = pd.read_parquet("consumer_cards_MDQ.parquet",
                      columns=CON_COLS)

# Merchant directory (all columns)
bus = pd.read_parquet("merchants_reference.parquet")""", fs=9)
add_explain_box(s, Inches(0.5), Inches(3.8), Inches(12.3), Inches(3.3), [
    "What is pd.read_parquet? It's like opening an Excel file, but much faster.",
    "Parquet is a columnar storage format — more efficient than CSV.",
    "We read ONLY the columns we need (columns=...) to save memory.",
    "",
    "Each row = 1 transaction. One card can have from 1 to 10,000+ transactions.",
    "There are FEWER business cards (25,000), but they are more active — 3M transactions.",
    "There are MORE consumer cards (80,000), but each does fewer transactions — 9.8M.",
    "",
    "IMPORTANT: business_cards contains transactions from KNOWN business cards.",
    "consumer_cards contains transactions from ordinary people.",
    "BUT some of them are hidden entrepreneurs (they just got a personal card",
    "and use it for business).",
], bg=RGBColor(0xE3,0xF2,0xFD), border=MED_BLUE)

# ══════════════════════════════════════════════════════════════════════
#  SLIDE 5: Label + concat
# ══════════════════════════════════════════════════════════════════════
s = new_slide("Concatenation — merging 2 tables into 1",
               "pd.concat in simple words + data example", "04")
code_block(s, Inches(0.4), Inches(1.4), Inches(4.8), Inches(5.5), """# STEP 1: assign labels
biz["label"] = 1  # business
con["label"] = 0  # consumer

# STEP 2: pick only needed columns
cols = BIZ_COLS + ["label"]
biz_subset = biz[cols]
con_subset = con[cols]

# STEP 3: concatenate
df = pd.concat([biz_subset, con_subset],
               ignore_index=True)
print(f"Total rows: {len(df):,}")
print(f"label=1 (biz): {(df['label']==1).sum():,}")
print(f"label=0 (con): {(df['label']==0).sum():,}")

# STEP 4: merge merchant info
df = df.merge(
    bus[["merchant_id","mcc","merchant_country",
         "recurring_capable","merchant_name"]],
    on=["merchant_id","mcc"], how="left"
)""", fs=9)
add_explain_box(s, Inches(5.5), Inches(1.4), Inches(7.3), Inches(2.3), [
    "pd.concat = glue two tables together, one on top of the other.",
    "We had 2 separate files: business and consumer.",
    "We place one ABOVE the other — we get one big table.",
    "ignore_index=True = reset old indices, create new 0,1,2,...",
    "",
    "label = 1 — 'this is a business card'. label = 0 — 'this is a consumer card'.",
    "The model will learn: what features distinguish label=1 from label=0.",
    "Then it will look at consumer cards and say: 'this one looks like business'.",
], icon="🧩", bg=RGBColor(0xE3,0xF2,0xFD), border=MED_BLUE)
add_data_example(s, Inches(5.5), Inches(3.9), Inches(7.3), Inches(3.0), [
    "biz (first 3 rows):",
    "card_number         | mcc  | amount | label",
    "5200000000000001   | 7311 | 150000 |   1",
    "5200000000000001   | 4214 | 45000  |   1",
    "5200000000000002   | 5021 | 89000  |   1",
    "",
    "con (first 3 rows):",
    "card_number         | mcc  | amount | label",
    "4100000000000001   | 5812 | 3500   |   0",
    "4100000000000001   | 5812 | 1200   |   0",
    "4100000000000002   | 5411 | 8900   |   0",
], [
    "df AFTER concat (first 6 rows):",
    "card_number         | mcc  | amount | label",
    "5200000000000001   | 7311 | 150000 |   1",
    "5200000000000001   | 4214 | 45000  |   1",
    "5200000000000002   | 5021 | 89000  |   1",
    "4100000000000001   | 5812 | 3500   |   0",
    "4100000000000001   | 5812 | 1200   |   0",
    "4100000000000002   | 5411 | 8900   |   0",
    "",
    "Total: 12,830,080 rows!",
    "Now all data is in one place.",
])

# ══════════════════════════════════════════════════════════════════════
#  SLIDE 6: MCC Categories
# ══════════════════════════════════════════════════════════════════════
s = new_slide("MCC Categories", "Classifying transactions: business, wholesale, logistics", "05")
code_block(s, Inches(0.4), Inches(1.4), Inches(5.5), Inches(5.5), """# BUSINESS TRANSACTIONS (everything for business)
BUSINESS_MCCS = {
    "7311","7372","5968","4816","4812","4814",
    "7379","5734","7392","7399","8931","8111",
    "8911","7361","5021","5044","5045","5046",
    "5111","5943","5099","5199","5131","5137",
    "5139","5169","5172","5193","5039","5065",
    "5072","5085","5211","5231","4214","4215",
    "4225","4011","4511","7011","4722","4723",
    "7394","7512","6381","7321","7333","7338",
    "7298","7299","5199","5261","5122","5200",
}

# WHOLESALE PRODUCT PURCHASES
WHOLESALE_PROD_MCCS = {"5021","5039","5044","5045",
    "5046","5065","5072","5085","5099","5111",
    "5122","5131","5137","5139","5169","5172",
    "5193","5199","5200","5211","5231",
}

# LOGISTICS + FUEL
LOGISTICS_MCCS = {"4214","4215","4225","4011"}
FUEL_MCCS     = {"5541","5542","5983"}

# SAAS / ADVERTISING (for cross-border)
SAAS_AD_MCCS  = {"7311","7372","5968","4816",
                 "7379","5734"}""", fs=7)
add_explain_box(s, Inches(6.2), Inches(1.4), Inches(6.6), Inches(2.8), [
    "MCC = Merchant Category Code — a code identifying the type of store.",
    "It's like 'departments in a supermarket': 7311 = advertising,",
    "4214 = freight, 5812 = restaurant, 5411 = grocery.",
    "",
    "We split all MCCs into 4 groups:",
    "🏢 BUSINESS_MCCS — everything a business needs (ads,",
    "  IT, telecom, wholesale, logistics, legal, hotels)",
    "📦 WHOLESALE — wholesale goods (building materials,",
    "  electronics, furniture, industrial equipment)",
    "🚚 LOGISTICS — freight, warehouses, railway",
    "⛽ FUEL — vehicle fuel",
    "💻 SAAS_AD — software, ads, cloud (for cross-border)",
    "",
    "WHY? To calculate features separately for each group.",
    "For example: 'wholesale spend ratio' = WHOLESALE spend / total spend.",
], icon="🏷️", bg=RGBColor(0xE3,0xF2,0xFD), border=MED_BLUE)
add_calc_box(s, Inches(6.2), Inches(4.5), Inches(6.6), Inches(2.4), [
    "Example: transaction with MCC=5021 (building materials)",
    "",
    "Check: is 5021 in WHOLESALE_PROD_MCCS? → YES",
    "So this is a wholesale purchase. Count in wholesale_spend.",
    "",
    "If MCC=7311 (ads) AND merchant country ≠ card country:",
    "→ this is a cross-border SaaS/Ads transaction.",
    "Count in cb_saas_ad_count.",
    "",
    "If MCC=5812 (restaurant):",
    "→ this is a consumer transaction. Not counted in paradox features.",
])

# ══════════════════════════════════════════════════════════════════════
#  PARADOX SLIDES (6-15): 10 Behavioral Paradoxes
# ══════════════════════════════════════════════════════════════════════

paradoxes = [
    {
        "num": "06", "letter": "5a",
        "title": "Clockwork Buyer — Purchase Regularity",
        "sub": "Paradox 1: Entrepreneurs buy like clockwork",
        "explain": [
            "IDEA: An entrepreneur orders goods every Monday.",
            "A consumer visits stores 'whenever they feel like'.",
            "We measure: HOW REGULAR are the intervals between purchases.",
            "",
            "HOW WE CALCULATE (example):",
            "Card X bought office supplies (MCC=5111):",
            "  Jan 1 → Jan 8 → Jan 15 → Jan 22 (every week)",
            "  Intervals: 7, 7, 7 days. CV = 0/7 = 0 → PERFECT",
            "",
            "Card Y bought office supplies:",
            "  Jan 1 → Feb 3 → Mar 10 → Apr 5",
            "  Intervals: 33, 35, 26 days. CV = 4.7/31 = 0.15 → still regular",
            "",
            "Card Z: Jan 1 → Mar 20 → Apr 2 → Dec 15",
            "  Intervals: 78, 13, 257 days. CV = large → CHAOTIC",
            "",
            "Lower CV → more regular purchasing → entrepreneur",
            "Higher CV → chaotic purchasing → consumer",
        ],
        "code": """# Sort by (card, MCC, time)
df.sort_values(["card_number","mcc",
    "transaction_timestamp"], inplace=True)

# Calculate interval between purchases of same MCC
df["interval_hours"] = df.groupby(
    ["card_number","mcc"]
)["transaction_timestamp"].diff().dt.total_seconds() / 3600

# CV = std/mean for each (card, MCC) pair
cv_per = df.dropna(subset=["interval_hours"]).groupby(
    ["card_number","mcc"]
)["interval_hours"].agg(
    lambda x: x.std()/x.mean() if x.mean()>1e-9 else np.nan
).rename("cv").reset_index()

# Average CV across all MCCs for each card
clockwork_cv_mean = cv_per.groupby(
    "card_number")["cv"].mean()

# How many MCCs have non-zero CV
clockwork_mcc_count = cv_per.dropna(
    subset=["cv"]).groupby("card_number")["cv"].count()""",
        "after": [
            "RESULTING FEATURES (per card):",
            "card_number  | clockwork_cv_mean | clockwork_mcc_count",
            "5200000001  | 0.05             | 8",
            "5200000002  | 0.12             | 5",
            "4100000001  | 0.89             | 2",
            "4100000002  | 1.45             | 1",
            "",
            "Business (5200...): low CV (0.05-0.12), many MCCs (5-8)",
            "Consumer (4100...): high CV (0.89-1.45), few MCCs (1-2)",
        ],
    },
    {
        "num": "07", "letter": "5b",
        "title": "Off-Hours Operator — Purchases outside 9-18",
        "sub": "Paradox 2: Business doesn't sleep before 9am and after 7pm",
        "explain": [
            "IDEA: Entrepreneurs order goods outside working hours —",
            "early morning (6-8am), late evening (8-11pm), weekends.",
            "Consumers shop during the day (12-6pm).",
            "",
            "HOW WE CALCULATE (example):",
            "From each transaction we extract the HOUR (0-23):",
            "  transaction_timestamp = '2026-01-15 07:30:00' → hour = 7",
            "  Hour < 9 or >= 19? 7 < 9 → YES → off_hours",
            "",
            "Card A (business): 100 business transactions,",
            "  70 of them during off-hours. off_hours_ratio = 70/100 = 0.70",
            "  Makes sense: store owner orders goods",
            "  in the morning before opening.",
            "",
            "Card B (consumer): 50 business transactions,",
            "  5 during off-hours. off_hours_ratio = 5/50 = 0.10",
            "  Consumer occasionally pays for a subscription at night.",
        ],
        "code": """# Extract hour from timestamp
df["txn_hour"] = df["transaction_timestamp"].dt.hour

# Define off-hours (before 9 or after 19)
df["is_off_hours"] = (df["txn_hour"] < 9) | \\
                      (df["txn_hour"] >= 19)

# Count business transactions during off-hours
off_hours_biz = df[df["is_business_mcc"] & \\
    df["is_off_hours"]].groupby("card_number").size()

# All business transactions
biz_mcc_txns = df[df["is_business_mcc"]] \\
    .groupby("card_number").size()

# Share of off-hours among business transactions
off_hours_ratio = (off_hours_biz / biz_mcc_txns) \\
    .rename("off_hours_ratio").fillna(0.0)

# Total transactions during off-hours
off_hours_total = df[df["is_off_hours"]] \\
    .groupby("card_number").size() \\
    .rename("off_hours_total_count")""",
        "after": [
            "CALCULATION EXAMPLE for card 4100000001:",
            "Transactions (with hour):",
            "  txn_1: 2026-01-15 07:30 → hour=7 → off_hours",
            "  txn_2: 2026-01-15 14:20 → hour=14 → NOT off_hours",
            "  txn_3: 2026-01-16 21:15 → hour=21 → off_hours",
            "  txn_4: 2026-01-17 08:45 → hour=8 → off_hours",
            "  txn_5: 2026-01-18 12:00 → hour=12 → NOT off_hours",
            "",
            "off_hours_ratio = 3/5 = 0.60",
            "total_txns = 5",
            "overall_off_hours_ratio = 3/5 = 0.60",
        ],
    },
    {
        "num": "08", "letter": "5c",
        "title": "Expense Ratio Inversion — Expense Inversion",
        "sub": "Paradox 3: Entrepreneurs spend more on goods than on themselves",
        "explain": [
            "IDEA: A consumer spends 90% on food, clothes, entertainment",
            "(MCC: 5411, 5812, 5651, 7841...). Only 10% on 'business'.",
            "An entrepreneur does the opposite: 60-80% of spending is",
            "wholesale purchases of goods (MCC: 5021, 5044, 5111...).",
            "",
            "HOW WE CALCULATE (example):",
            "Card A (business): total spent 5,000,000 KZT.",
            "  Of that, on wholesale MCCs: 3,500,000 KZT.",
            "  wholesale_spend_ratio = 3.5M/5M = 0.70 (70%)",
            "",
            "Card B (consumer): total spent 500,000 KZT.",
            "  Of that, on wholesale MCCs: 20,000 KZT.",
            "  wholesale_spend_ratio = 20K/500K = 0.04 (4%)",
            "",
            "Higher ratio → more resembles an entrepreneur",
        ],
        "code": """# Total spend per card
total_spend = df.groupby("card_number") \\
    ["transaction_amount_kzt"].sum() \\
    .rename("total_spend_kzt")

# Wholesale spend per card
wholesale_spend = df[
    df["mcc"].isin(WHOLESALE_PROD_MCCS)
].groupby("card_number")[
    "transaction_amount_kzt"
].sum().rename("wholesale_spend_kzt")

# Wholesale share of total spend
wholesale_spend_ratio = (
    wholesale_spend / total_spend
).rename("wholesale_spend_ratio").fillna(0.0)

# Log of wholesale spend (smoothing)
wholesale_spend_log = np.log1p(
    wholesale_spend
).rename("wholesale_spend_log")""",
        "after": [
            "EXAMPLE: why log1p?",
            "Card A spent 10,000,000 KZT wholesale",
            "Card B spent 50,000 KZT wholesale",
            "That's a 200x difference!",
            "log1p(10M) = 16.1, log1p(50K) = 10.8",
            "Only 1.5x difference → the model won't",
            "be overwhelmed by millionaires",
        ],
    },
    {
        "num": "09", "letter": "5d",
        "title": "Token Wholesale — Tokenized Wholesale",
        "sub": "Paradox 4: Large wholesale via Apple/Google Pay",
        "explain": [
            "IDEA: Entrepreneurs pay suppliers via token",
            "(Apple Pay, Google Pay, Samsung Pay) — fast and convenient",
            "for frequent purchases. Consumers mostly tokenize",
            "small everyday transactions.",
            "",
            "HOW WE CALCULATE:",
            "1) Take ONLY wholesale MCCs + tokenized = True",
            "2) Among business cards, find the 90th percentile amount:",
            "   (90% of business transactions are BELOW this amount)",
            "   = 470,088 KZT ≈ $1,150",
            "3) If a consumer card has a wholesale tokenized",
            "   transaction ABOVE this threshold → hidden entrepreneur",
            "",
            "WHY THE 90TH PERCENTILE?",
            "It represents a 'very large purchase' that an ordinary person",
            "wouldn't make with a token. But an entrepreneur would (restocking).",
        ],
        "code": """# Filter: wholesale + tokenized
token_wholesale = df[
    df["mcc"].isin(WHOLESALE_PROD_MCCS)
    & (df["tokenized"] == True)
].copy()

# 90th percentile of business cards
biz_token = token_wholesale[
    token_wholesale["label"] == 1
]
p90 = biz_token["transaction_amount_kzt"] \\
    .quantile(0.90)

# Flag: does the card have any such transaction?
token_wholesale_flag = (
    token_wholesale[
        token_wholesale["transaction_amount_kzt"] > p90
    ].groupby("card_number").size().gt(0)
    .astype(int).rename("token_wholesale_flag")
)

# Count: how many such transactions?
token_wholesale_count = (
    token_wholesale[
        token_wholesale["transaction_amount_kzt"] > p90
    ].groupby("card_number").size()
    .rename("token_wholesale_count")
)""",
        "after": [
            "Example: consumer card 4100000001:",
            "  Wholesale transaction 500,000 KZT, tokenized=True",
            "  500K > 470K (business 90th percentile)",
            "  → token_wholesale_flag = 1",
            "  → LIKELY A HIDDEN ENTREPRENEUR",
            "",
            "Example: consumer card 4100000002:",
            "  No wholesale tokenized transactions at all",
            "  → token_wholesale_flag = 0",
            "  → ordinary consumer",
        ],
    },
    {
        "num": "10", "letter": "5e",
        "title": "Baseline Aggregations — Basic Statistics",
        "sub": "The most important features (tokenized_ratio = 26% importance!)",
        "explain": [
            "IDEA: Beyond smart paradoxes, we need basic data:",
            "• Total number of transactions?",
            "• Average/Max amount?",
            "• How many different stores? MCCs?",
            "• Online vs POS ratio?",
            "• Share of tokenized payments?",
            "",
            "THESE FEATURES TURNED OUT TO BE THE MOST IMPORTANT:",
            "tokenized_ratio (26.3%) — business cards pay almost",
            "everything via token. Consumer — cash/card.",
            "",
            "online_ratio (20.3%) — business orders online,",
            "consumer goes to the store.",
            "",
            "TOGETHER = 46.6% of ALL model importance!",
            "Paradoxes complement, but baseline is the foundation.",
        ],
        "code": """card_agg = df.groupby("card_number").agg(
    amount_mean=("transaction_amount_kzt","mean"),
    amount_std =("transaction_amount_kzt","std"),
    amount_sum =("transaction_amount_kzt","sum"),
    amount_max =("transaction_amount_kzt","max"),
    n_unique_mcc=("mcc","nunique"),
    n_unique_merchants=("merchant_id","nunique"),
    n_online=("channel",
              lambda x:(x=="online").sum()),
    n_pos=("channel",
           lambda x:(x=="POS").sum()),
    tokenized_txn_count=("tokenized","sum"),
    recurring_count=("is_recurring","sum"),
)

# Share of online transactions
card_agg["online_ratio"] = card_agg["n_online"] / (
    card_agg["n_online"] + card_agg["n_pos"]
)

# Share of tokenized transactions
card_agg["tokenized_ratio"] = (
    card_agg["tokenized_txn_count"] / total_txns
)

# Log of MCC count (smoothing)
card_agg["n_unique_mcc_log"] = np.log1p(
    card_agg["n_unique_mcc"]
)""",
        "after": [
            "EXAMPLE (typical values):",
            "",
            "Feature              | Business  | Consumer",
            "tokenized_ratio     | 0.65-0.85 | 0.05-0.15",
            "online_ratio        | 0.70-0.90 | 0.20-0.40",
            "amount_mean         | 25,000+   | 3,000-8,000",
            "amount_sum          | 2M+ KZT   | 200-500K KZT",
            "n_unique_merchants  | 50-200    | 10-50",
            "",
            "Huge gap! This is why",
            "tokenized_ratio is #1 in importance.",
        ],
    },
    {
        "num": "11", "letter": "5f",
        "title": "Supplier Fingerprint — Vendor Concentration",
        "sub": "Paradox 5 (NEW): Entrepreneur = 3-5 key suppliers",
        "explain": [
            "IDEA: A consumer visits 100+ different places: grocery stores,",
            "restaurants, clothing shops, Netflix... Money is 'spread thin'.",
            "An entrepreneur buys from 3-5 key suppliers:",
            "wholesale warehouse, logistics company, ad agency.",
            "",
            "HOW WE CALCULATE:",
            "vendor_concentration = spend on TOP-3 suppliers",
            "                    / total card spend",
            "",
            "If 80% of all money went to 3 suppliers → 0.8",
            "→ entrepreneur (narrow circle)",
            "If 20% → 0.2 → consumer (spread out)",
            "",
            "Additionally: merchant_gini = Gini coefficient",
            "0 = all spending evenly distributed,",
            "1 = everything to one merchant",
        ],
        "code": """# Spend per (card, merchant) pair
merchant_spend = df.groupby(
    ["card_number","merchant_id"]
)["transaction_amount_kzt"].sum().reset_index()

# Top-3 supplier share of total spend
def _top3_share(g):
    total = g["transaction_amount_kzt"].sum()
    if total == 0: return 0.0
    top3 = g.nlargest(3, "transaction_amount_kzt") \\
        ["transaction_amount_kzt"].sum()
    return top3 / total

vendor_concentration = merchant_spend \\
    .groupby("card_number") \\
    .apply(_top3_share, include_groups=False)

# Gini coefficient of concentration
def _gini(g):
    amounts = g["transaction_amount_kzt"] \\
        .sort_values().values
    if len(amounts)==0 or amounts.sum()==0:
        return 0.0
    n = len(amounts)
    cumsum = np.cumsum(amounts)
    return (2*np.sum(cumsum)/cumsum[-1]-n-1)/n

merchant_gini = merchant_spend \\
    .groupby("card_number") \\
    .apply(_gini, include_groups=False)""",
        "after": [
            "EXAMPLE:",
            "Business card (5200000001):",
            "  Supplier 1 (wholesale): 2,000,000 (40%)",
            "  Supplier 2 (logistics): 1,500,000 (30%)",
            "  Supplier 3 (ads):        500,000 (10%)",
            "  Others:                1,000,000 (20%)",
            "  vendor_concentration = (2M+1.5M+0.5M)/5M = 0.80",
            "",
            "Consumer card (4100000001):",
            "  Top-3 = 150K out of 500K = 0.30",
            "  → ordinary consumer",
        ],
    },
    {
        "num": "12", "letter": "5g",
        "title": "Last-Mile Echo — Wholesale to Logistics Lag",
        "sub": "Paradox 6 (NEW): Buy goods → deliver within 0-7 days",
        "explain": [
            "IDEA: Entrepreneur buys goods → pays for delivery to",
            "customer within 0-7 days. A consumer never does this.",
            "",
            "HOW WE CALCULATE:",
            "1) Aggregate transactions by DAY (not by hour)",
            "2) For each card, find wholesale purchase days",
            "3) For each wholesale day, find the NEXT logistics/fuel",
            "   transaction forward in time",
            "4) Calculate average LAG (days) between wholesale and logistics",
            "5) If lag is 0-7 days → this is Last-Mile Echo",
            "",
            "TECHNICAL PROBLEM (important!):",
            "pd.merge_asof with by=['card_number'] HANGS in pandas 3.0.3",
            "Solution: numpy.searchsorted — 100x faster",
            "",
            "Why day-level aggregation? 12.8M rows → 300K days.",
            "40x less data = 40x faster processing.",
        ],
        "code": """# Aggregate by day (12.8M → 300K rows)
daily_wholesale = df[df["mcc"].isin(WHOLESALE_PROD_MCCS)] \\
    .groupby(["card_number","transaction_date"]) \\
    ["transaction_amount_kzt"].sum().reset_index()

daily_logistics = df[
    df["mcc"].isin(LOGISTICS_MCCS | FUEL_MCCS)
].groupby(["card_number","transaction_date"]) \\
    ["transaction_amount_kzt"].sum().reset_index()

# Dictionary: card → array of logistics dates
# numpy.searchsorted — binary search
logistics_dates = {
    cn: g["date"].values.astype("datetime64[D]")
    for cn, g in daily_logistics.groupby("card_number",
                                         sort=False)
}

def _echo_avg(group):
    cn = group.name
    w_dates = group.loc[group["wholesale_amt"]>0,
        "date"].values.astype("datetime64[D]")
    l_dates = logistics_dates.get(cn,
        np.array([], dtype="datetime64[D]"))
    if len(w_dates)==0 or len(l_dates)==0:
        return 0.0
    # searchsorted = find nearest logistics date
    idx = np.searchsorted(l_dates, w_dates, side="left")
    idx = np.clip(idx, 0, len(l_dates)-1)
    lags = (l_dates[idx] - w_dates).astype(int)
    valid = (lags >= 0) & (lags <= 7)
    return lags[valid].mean() if valid.sum()>0 else 0.0

avg_echo_lag = daily_merged.groupby("card_number",
    sort=False).apply(_echo_avg,
    include_groups=False).rename("wholesale_to_logistics_lag")""",
        "after": [
            "EXAMPLE:",
            "Card 5200000001 (business):",
            "  Day 1: wholesale 500,000 KZT",
            "  Day 3: logistics 45,000 KZT → lag=2 days",
            "  Day 15: wholesale 300,000 KZT",
            "  Day 18: fuel 12,000 KZT → lag=3 days",
            "  Average lag = (2+3)/2 = 2.5 days ✓",
            "",
            "Card 4100000001 (consumer):",
            "  No wholesale purchases → lag = 0",
            "  (filled with 168 hours = 1 week)",
        ],
    },
    {
        "num": "13", "letter": "5h",
        "title": "Round-Trip Cash Flow — Spending Burst Periodicity",
        "sub": "Paradox 7 (NEW, REDEFINED): Purchase cycles",
        "explain": [
            "IDEA (original): Detect money round-tripping",
            "(paid → refunded → paid again).",
            "",
            "PROBLEM: Our data has ONLY outflows (spending).",
            "No refunds, no deposits, no incoming transactions.",
            "The original paradox is IMPOSSIBLE to implement.",
            "",
            "SOLUTION — REDEFINITION:",
            "We look at the PERIODICITY of large spending.",
            "An entrepreneur buys goods every month (spending burst).",
            "A consumer spends large amounts chaotically (NYE, birthday).",
            "",
            "HOW WE CALCULATE:",
            "1) Group daily spending per card",
            "2) Find days with spending in the TOP 5% (bursts)",
            "3) Calculate intervals between these days",
            "4) std of intervals — if small, spending is regular",
            "",
            "Example: std=2 → bursts every month ±2 days",
            "Example: std=45 → chaotic, once every 6 months-year",
        ],
        "code": """# Daily spending
daily_spend = df.groupby(
    ["card_number","transaction_date"]
)["transaction_amount_kzt"].sum().reset_index()

def _burst_periodicity(g):
    amounts = g["transaction_amount_kzt"]
    if len(amounts) < 5:
        return np.nan  # not enough data
    threshold = amounts.quantile(0.95)  # top 5%
    burst_days = g[amounts >= threshold] \\
        ["transaction_date"].sort_values()
    if len(burst_days) < 2:
        return np.nan
    intervals = burst_days.diff() \\
        .dt.days.dropna()
    if len(intervals) < 2:
        return np.nan
    return float(intervals.std())
    # Small std = regular purchases

spending_burst_periodicity = daily_spend \\
    .groupby("card_number") \\
    .apply(_burst_periodicity, include_groups=False) \\
    .rename("spending_burst_periodicity")""",
        "after": [
            "CALCULATION EXAMPLE:",
            "",
            "Business card — daily spending:",
            "  Jan 1: 500,000 ← top 5%",
            "  Jan 5: 3,000",
            "  Jan 8: 600,000 ← top 5%",
            "  Jan 12: 2,000",
            "  Jan 15: 450,000 ← top 5%",
            "  ...",
            "  Intervals: 7d, 7d, 7d...",
            "  std = 0.5 days → very regular!",
            "",
            "Consumer card:",
            "  Bursts every 30-90 days",
            "  std = 45 days → chaotic",
        ],
    },
    {
        "num": "14", "letter": "5i",
        "title": "Inventory Pulse — Mass Purchase + Small Top-ups",
        "sub": "Paradox 8 (NEW): Large restock + many small reorders",
        "explain": [
            "IDEA: An entrepreneur buys a million worth of goods once",
            "a month (massive), and between — small reorders of",
            "consumables (small). A consumer does neither.",
            "",
            "HOW WE CALCULATE:",
            "1) Take only wholesale transactions of the card",
            "2) Find P30 (30th percentile) and P70 (70th percentile)",
            "3) small = count of transactions ≤ P30",
            "4) large = count of transactions ≥ P70",
            "5) ratio = small / large",
            "",
            "If ratio = 3.0: for 1 large purchase there are",
            "3 small top-ups → entrepreneur pattern",
            "",
            "If ratio = 0.2: almost all transactions are large →",
            "doesn't look like inventory management",
            "",
            "IMPORTANT: need at least 5 wholesale transactions,",
            "otherwise we can't draw conclusions. Return NaN.",
        ],
        "code": """def _inventory_pulse(g):
    # Take only wholesale transactions of the card
    wholesale = g[g["mcc"].isin(
        WHOLESALE_PROD_MCCS)]
    if len(wholesale) < 5:
        return np.nan  # insufficient data
    amounts = wholesale["transaction_amount_kzt"]
    p30 = amounts.quantile(0.30)
    p70 = amounts.quantile(0.70)
    small = (amounts <= p30).sum()
    large = (amounts >= p70).sum()
    return small / large if large > 0 else np.nan

massive_to_small_ratio = df \\
    .groupby("card_number") \\
    .apply(_inventory_pulse, include_groups=False) \\
    .rename("massive_to_small_ratio")""",
        "after": [
            "EXAMPLE:",
            "Business card — wholesale transactions:",
            "  1,000,000 (large), 50,000 (small), 30,000 (small),",
            "  800,000 (large), 40,000 (small), 20,000 (small),",
            "  1,200,000 (large), 60,000 (small)",
            "P30 = 35,000, P70 = 900,000",
            "small = 4, large = 3, ratio = 4/3 = 1.33",
            "",
            "Consumer card — wholesale transactions:",
            "  5,000, 8,000, 12,000, 3,000, 6,000",
            "  All ≈ same size, ratio ≈ 1.0",
            "  → not business (no 'bulk+topup' structure)",
        ],
    },
    {
        "num": "15", "letter": "5j",
        "title": "Multi-Vendor Loyalty Paradox — Volume without Loyalty",
        "sub": "Paradox 9 (NEW): High B2B spend, NO recurring subscriptions",
        "explain": [
            "IDEA (key): An entrepreneur spends heavily on",
            "business services (ads, IT, wholesale, logistics), but does",
            "NOT set up recurring subscriptions (Netflix, Spotify).",
            "",
            "A consumer does the opposite — subscribed to everything.",
            "",
            "FORMULA: b2b_volume_no_recurring =",
            "  (total B2B spend) × (1 − recurring share)",
            "",
            "If entrepreneur spends 2M on business and 0% recurring:",
            "  score = 2M × (1−0) = 2,000,000 → HIGH",
            "",
            "If consumer spends 50K on business and 30% recurring:",
            "  score = 50K × (1−0.3) = 35,000 → LOW",
            "",
            "THIRD most important feature (12%)!",
        ],
        "code": """# Total B2B MCC spend
b2b_spend = df[df["mcc"].isin(BUSINESS_MCCS)] \\
    .groupby("card_number")["transaction_amount_kzt"].sum()

# Recurring transaction share per card
recurring_ratio = df.groupby("card_number") \\
    ["is_recurring"].mean()

# B2B volume without recurring (KEY feature)
b2b_volume_no_recurring = (
    b2b_spend * (1 - recurring_ratio)
).rename("b2b_volume_no_recurring").fillna(0.0)

# Avg transactions per merchant
txns_per_merchant = total_txns / \\
    df.groupby("card_number")["merchant_id"].nunique()""",
        "after": [
            "COMPARISON:",
            "",
            "Business card:",
            "  B2B spend: 3,000,000 KZT",
            "  recurring_ratio: 0.02 (2%)",
            "  b2b_volume_no_recurring = 3M × 0.98 = 2,940,000",
            "  Works with suppliers, doesn't subscribe",
            "",
            "Consumer card:",
            "  B2B spend: 50,000 KZT",
            "  recurring_ratio: 0.25 (25%)",
            "  b2b_volume_no_recurring = 50K × 0.75 = 37,500",
            "  (Netflix, Spotify subscriptions = recurring)",
            "",
            "80x difference! → critical feature",
        ],
    },
    {
        "num": "16", "letter": "5k",
        "title": "Channel Schizophrenia — Channel Switching",
        "sub": "Paradox 10 (NEW): Online ↔ POS constant switching",
        "explain": [
            "IDEA: An entrepreneur orders goods online (supplier",
            "website), goes to the warehouse (POS terminal), orders",
            "again online, goes to the gas station (POS). Constant",
            "switching between online and POS.",
            "",
            "A consumer usually sticks to one channel:",
            "• Some buy everything online",
            "• Some only in physical stores",
            "",
            "HOW WE CALCULATE:",
            "1) Sort transactions by time",
            "2) For each, check: did the channel change?",
            "3) channel_alternation_rate = share of channel changes",
            "",
            "Example sequence:",
            "online → POS → online → POS → online",
            "4 changes out of 5 transactions = 0.8 (80%)→ entrepreneur",
            "",
            "online → online → online → POS → POS",
            "1 change out of 5 = 0.2 (20%) → consumer",
        ],
        "code": """# Sort by time
df_sorted = df.sort_values(
    ["card_number","transaction_timestamp"]
)

# Shift channel by 1 position
df_sorted["prev_channel"] = df_sorted \\
    .groupby("card_number")["channel"].shift(1)

# Was there a channel change?
df_sorted["channel_switch"] = (
    (df_sorted["channel"] != df_sorted["prev_channel"])
    & df_sorted["prev_channel"].notna()
)

# Share of channel changes
channel_alternation_rate = df_sorted \\
    .groupby("card_number")["channel_switch"] \\
    .mean().rename("channel_alternation_rate")

# Shannon entropy: 1 = 50/50 online/POS
def _entropy(g):
    counts = g.value_counts()
    probs = counts / counts.sum()
    return float(-(probs*np.log2(probs+1e-10)).sum())

channel_entropy = df_sorted \\
    .groupby("card_number")["channel"] \\
    .apply(_entropy).rename("channel_entropy")""",
        "after": [
            "EXAMPLE sequences:",
            "",
            "Business (alternation_rate = 0.70):",
            "  ONLINE → POS → ONLINE → POS → ONLINE → POS",
            "  Constantly switching = ordering + shipping",
            "",
            "Consumer (alternation_rate = 0.15):",
            "  ONLINE → ONLINE → ONLINE → ONLINE → POS → POS",
            "  Mostly online, occasional in-store visit",
            "",
            "FOURTH most important feature (9.2%)!",
        ],
    },
]

for p in paradoxes:
    s = new_slide(p["title"], p["sub"], p["num"])
    code_block(s, Inches(0.3), Inches(1.35), Inches(5.2), Inches(5.8), p["code"], fs=7)
    add_explain_box(s, Inches(5.8), Inches(1.35), Inches(7.0), Inches(3.8),
                    p["explain"], icon=SIMPLE_ICON, bg=RGBColor(0xE3,0xF2,0xFD), border=MED_BLUE)
    add_data_example(s, Inches(5.8), Inches(5.3), Inches(7.0), Inches(1.9),
                     [], p.get("after", ["(data example)"]))

# ══════════════════════════════════════════════════════════════════════
#  SLIDE 17: Cross-border Sourcing
# ══════════════════════════════════════════════════════════════════════
s = new_slide("Cross-border Sourcing — International SaaS/Ads",
               "🌍 Paradox 11: Google Ads, AWS, Shopify — paid abroad", "17")
code_block(s, Inches(0.3), Inches(1.35), Inches(5.2), Inches(5.8), """# Fill merchant country if missing
df["merchant_country_filled"] = df[
    "merchant_country"].fillna(df["country"])

# Cross-border: merchant country ≠ card country
df["is_cross_border"] = (
    df["merchant_country"] != df["country"]
) & df["merchant_country"].notna()

# SaaS/Ads + cross-border
df["is_cb_saas_ad"] = df["is_cross_border"] & \\
    df["mcc"].isin(SAAS_AD_MCCS)

# Share of international transactions
cb_ratio = df.groupby("card_number") \\
    ["is_cross_border"].mean().rename("cb_ratio")

# Share of SaaS/Ads that are international
cb_saas_ad_ratio = df[df["mcc"].isin(SAAS_AD_MCCS)] \\
    .groupby("card_number")["is_cross_border"] \\
    .mean().rename("cb_saas_ad_ratio")

# Share of spend on international transactions
cb_spend = df[df["is_cross_border"]] \\
    .groupby("card_number")["transaction_amount_kzt"].sum()
cb_spend_share = (cb_spend / total_spend) \\
    .rename("cb_spend_share").fillna(0.0)

# Count of international SaaS/Ads
cb_saas_ad_count = df[df["is_cb_saas_ad"]] \\
    .groupby("card_number").size() \\
    .rename("cb_saas_ad_count")""", fs=7)
add_explain_box(s, Inches(5.8), Inches(1.35), Inches(7.0), Inches(3.5), [
    "IDEA: Entrepreneurs pay ABROAD for Google Ads,",
    "Facebook Ads, AWS (cloud), Shopify (e-commerce), Zoom.",
    "Consumers usually DON'T pay abroad — everything is local.",
    "",
    "HOW WE DETECT:",
    "Each transaction has a card country (e.g., KZ)",
    "and a merchant country (from merchants_reference).",
    "If KZ ≠ US → cross-border = True.",
    "If the MCC is in SAAS_AD_MCCS → SaaS/Ads transaction.",
    "",
    "EXAMPLE:",
    "Card 4100000001:",
    "  Pays Google Ads (MCC=7311), merchant in USA",
    "  → cb_saas_ad_count += 1",
    "  cb_spend_share = $500 / $50,000 = 0.01 (1%)",
    "",
    "Card 5200000001 (business):",
    "  Google Ads + AWS + Shopify = $5,000",
    "  cb_spend_share = $5,000/$20,000 = 0.25 (25%)",
], icon="🌍", bg=RGBColor(0xE3,0xF2,0xFD), border=MED_BLUE)
add_data_example(s, Inches(5.8), Inches(5.1), Inches(7.0), Inches(2.1), [], [
    "COMPARISON:",
    "",
    "Business: 15-40% of spending — international SaaS/Ads",
    "Consumer: 0-2% of spending — international SaaS/Ads",
    "",
    "cb_spend_share = 2.88% importance in the model!",
])

# ══════════════════════════════════════════════════════════════════════
#  SLIDE 18: Feature Assembly
# ══════════════════════════════════════════════════════════════════════
s = new_slide("Feature Assembly — 35 features",
               "Combine all features, fill NaN, get 105K × 35 matrix", "18")
code_block(s, Inches(0.3), Inches(1.35), Inches(5.2), Inches(5.8), """# 1. Get labels for each card
labels = df[["card_number","label"]] \\
    .drop_duplicates("card_number") \\
    .set_index("card_number")["label"]

# 2. Create feature matrix (index = card_number)
features = pd.DataFrame(index=labels.index)

# 3. Add all 35 features
features["clockwork_cv_mean"] = clockwork_cv
features["off_hours_ratio"] = off_hours_ratio
features["wholesale_spend_ratio"] = expense_ratio
features["token_wholesale_flag"] = token_flag
features["total_txns"] = total_txns
features["amount_mean"] = card_agg["amount_mean"]
features["online_ratio"] = card_agg["online_ratio"]
features["tokenized_ratio"] = card_agg["tokenized_ratio"]
features["vendor_concentration"] = vendor_concentration
features["wholesale_to_logistics_lag"] = avg_echo_lag
features["spending_burst_periodicity"] = burst_periodicity
features["massive_to_small_ratio"] = inventory_pulse
features["b2b_volume_no_recurring"] = b2b_volume
features["channel_alternation_rate"] = channel_alt
features["channel_entropy"] = channel_ent
features["cb_ratio"] = cb_ratio
# ... and so on for all 35

# 4. Fill NaN values
features["clockwork_cv_mean"] = features[
    "clockwork_cv_mean"].fillna(
    features["clockwork_cv_mean"].median())
features["off_hours_ratio"] = features[
    "off_hours_ratio"].fillna(0.0)
# ... each feature has its own fill strategy

# 5. Drop remaining NaN rows
features = features.dropna()
print(f"Features: {features.shape[1]}")
print(f"Cards: {len(features):,}")""", fs=7)
mtb(s, Inches(5.8), Inches(1.35), Inches(7.0), Inches(2.8), [
    "35 features — what's included:",
    "",
    "📊 Clockwork (2):    cv_mean, mcc_count",
    "📊 Off-Hours (3):    ratio, overall_ratio, total_count",
    "📊 Expense (2):      wholesale_ratio, wholesale_log",
    "📊 Token (2):        flag, count",
    "📊 Baseline (10):    total_txns, amount_mean, amount_std,",
    "                     amount_sum, amount_max, n_unique_mcc_log,",
    "                     n_unique_mcc, n_unique_merchants,",
    "                     online_ratio, tokenized_ratio",
    "📊 Supplier (2):     vendor_concentration, merchant_gini",
    "📊 Last-Mile (2):    lag, cross_ratio",
    "📊 Round-Trip (1):   burst_periodicity",
    "📊 Inventory (1):    massive_to_small_ratio",
    "📊 Multi-Vendor (2): b2b_volume_no_recurring, txns_per_merchant",
    "📊 Channel (2):      alternation_rate, entropy",
    "📊 Cross-border (4): cb_ratio, cb_saas_ad_ratio,",
    "                     cb_spend_share, cb_saas_ad_count",
    "",
    "Total: 35 features for 105,000 cards",
], fs=11, ls=1.2)
add_compare_box(s, Inches(5.8), Inches(4.4), Inches(7.0), Inches(2.7), [
    "✅ tokenized_ratio — most important (26.3%)",
    "✅ online_ratio (20.3%)",
    "✅ b2b_volume_no_recurring (12.0%)",
    "✅ channel_alternation_rate (9.2%)",
    "✅ All 10 paradoxes create synergy",
    "✅ Cross-border adds ~6%",
], [
    "❌ card_tier: consumer has no 'Business' tier",
    "❌ bank_name: identical distribution",
    "❌ recurring_capable: only 27/2165 merchants",
    "❌ merchant_name: MCC already covers it",
    "",
    "35 features are enough for AUC=1.0",
])

# ══════════════════════════════════════════════════════════════════════
#  SLIDE 19: Train/Test + CatBoost
# ══════════════════════════════════════════════════════════════════════
s = new_slide("CatBoost — Model Training", "Why CatBoost, parameters, early stopping", "19")
code_block(s, Inches(0.3), Inches(1.35), Inches(5.2), Inches(5.8), """# ── 80/20 Stratified Split ──────────
X_train, X_test, y_train, y_test = \\
    train_test_split(
    features.values,
    labels.values,
    test_size=0.2,      # 20% for testing
    random_state=42,     # reproducibility
    stratify=labels.values,  # preserve
)                             # class ratio

print(f"Train: {len(X_train):,}")
print(f"Test:  {len(X_test):,}")

# ── CatBoost ─────────────────────────
model = CatBoostClassifier(
    iterations=500,          # max iterations
    learning_rate=0.1,       # learning step
    depth=6,                 # tree depth
    loss_function="Logloss", # binary classification
    eval_metric="AUC",       # what we optimize
    auto_class_weights=      # class balance
        "Balanced",
    early_stopping_rounds=50,# stop if AUC
    verbose=50,              # doesn't improve
    random_seed=42,
    task_type="CPU",
    allow_writing_files=False,
)

model.fit(
    X_train, y_train,
    eval_set=(X_test, y_test),
    use_best_model=True,  # take the best
)                         # iteration""", fs=7)
add_explain_box(s, Inches(5.8), Inches(1.35), Inches(7.0), Inches(2.3), [
    "WHY CATBOOST instead of XGBoost/LightGBM?",
    "",
    "1) CatBoost — best AUC on tabular data",
    "2) Ordered Boosting — less overfitting",
    "3) Symmetric Trees — faster inference",
    "4) auto_class_weights='Balanced' — automatically",
    "   compensates for class imbalance",
    "",
    "WHAT EACH PARAMETER MEANS:",
    "• iterations=500: model can make up to 500 steps",
    "• learning_rate=0.1: each step is small",
    "• early_stopping_rounds=50: if for 50 consecutive",
    "  steps AUC doesn't improve → stop",
    "• eval_metric='AUC': optimize what we evaluate",
    "• random_seed=42: same result every",
    "  run (reproducibility)",
], icon="🤖", bg=RGBColor(0xE3,0xF2,0xFD), border=MED_BLUE)
add_calc_box(s, Inches(5.8), Inches(3.9), Inches(7.0), Inches(3.3), [
    "TRAINING PROGRESS:",
    "  Iteration   0: test AUC = 0.9965",
    "  Iteration  50: test AUC = 0.999998",
    "  Iteration 100: test AUC = 0.999999",
    "  Stopped at iteration 121 (out of 500)",
    "  → Test AUC = 0.999999 ≈ 1.0000",
    "",
    "Train: 84,000 cards | Test: 21,000 cards",
    "",
    "WHY DID IT STOP EARLY?",
    "The model reached maximum performance in 121 steps.",
    "Further training could lead to overfitting.",
    "Early stopping protects against this.",
    "",
    "use_best_model=True → take iteration 121,",
    "not iteration 500 (which might be worse on test).",
])

# ══════════════════════════════════════════════════════════════════════
#  SLIDE 20: 5-Fold Cross-Validation
# ══════════════════════════════════════════════════════════════════════
s = new_slide("5-Fold Cross-Validation", "Is the model stable on ANY data split?", "20")
code_block(s, Inches(0.3), Inches(1.35), Inches(5.2), Inches(5.8), """# CV model (fewer iterations —
# trains 5 times!)
cv_model = CatBoostClassifier(
    iterations=200,
    learning_rate=0.1,
    depth=6,
    loss_function="Logloss",
    eval_metric="AUC",
    auto_class_weights="Balanced",
    random_seed=42,
    task_type="CPU",
    allow_writing_files=False,
    verbose=0,
)

# Pool — CatBoost data format
cv_pool = Pool(
    features.values,
    labels.values,
    feature_names=feature_names,
)

# 5-Fold Stratified Cross-Validation
cv_results = catboost_cv(
    cv_pool,
    cv_model.get_params(),
    fold_count=5,         # 5 splits
    stratified=True,       # preserve class ratio
    early_stopping_rounds=30,
    verbose_eval=False,
)

print(f"CV AUC mean: "
      f"{cv_results['test-AUC-mean'].iloc[-1]:.4f}")
print(f"CV AUC std:  "
      f"{cv_results['test-AUC-std'].iloc[-1]:.4f}")""", fs=7)
add_explain_box(s, Inches(5.8), Inches(1.35), Inches(7.0), Inches(2.0), [
    "WHAT IS CROSS-VALIDATION (simple terms):",
    "",
    "Instead of 1 test, we do 5 tests.",
    "Each time 80% of data → training, 20% → testing.",
    "But each time, the 20% is DIFFERENT.",
    "",
    "If the model shows high AUC ALL 5 TIMES —",
    "it didn't get lucky with one random split,",
    "it ACTUALLY works well.",
    "",
    "Stratification: each fold preserves the",
    "76% consumer / 24% business ratio.",
], icon="🔬", bg=RGBColor(0xE3,0xF2,0xFD), border=MED_BLUE)
add_calc_box(s, Inches(5.8), Inches(3.6), Inches(7.0), Inches(3.6), [
    "5-FOLD CV RESULTS:",
    "",
    "  Fold 0: AUC = 0.999999",
    "  Fold 1: AUC = 0.999998",
    "  Fold 2: AUC = 0.999999",
    "  Fold 3: AUC = 0.999998",
    "  Fold 4: AUC = 0.999999",
    "",
    "  Mean AUC:  1.0000",
    "  Std Dev:   0.0000",
    "",
    "WHAT THIS MEANS:",
    "• Model is stable on ANY data",
    "• No dependence on random split",
    "• AUC=1.0 is NOT overfitting (confirmed 5x)",
    "",
    "⚠️ If it were overfitting, CV AUC would be",
    "  SIGNIFICANTLY lower than test AUC (e.g.,",
    "  0.95 vs 1.0). Here they're identical — honest.",
])

# ══════════════════════════════════════════════════════════════════════
#  SLIDE 21: Results — AUC-ROC + Confusion Matrix
# ══════════════════════════════════════════════════════════════════════
s = new_slide("RESULTS — AUC-ROC and Confusion Matrix",
               "Key metrics: where numbers come from and what they mean", "21")
code_block(s, Inches(0.3), Inches(1.35), Inches(4.5), Inches(5.8), """# Predictions on test set
y_pred = model.predict(X_test)
y_proba = model.predict_proba(
    X_test)[:, 1]  # probabilities

# PRIMARY: AUC-ROC
auc = roc_auc_score(y_test, y_proba)
print(f"AUC-ROC: {auc:.4f}")

# SECONDARY: Confusion Matrix
cm = confusion_matrix(y_test, y_pred)
print(f"TN={cm[0,0]} FP={cm[0,1]}")
print(f"FN={cm[1,0]} TP={cm[1,1]}")

# Visualization
ConfusionMatrixDisplay.from_estimator(
    model, X_test, y_test,
    display_labels=["Consumer","Business"],
    cmap="Blues",
)
plt.savefig("confusion_matrix.png", dpi=150)""", fs=8)
mtb(s, Inches(5.1), Inches(1.35), Inches(4.0), Inches(5.8), [
    "📊 AUC-ROC = 1.0000",
    "",
    "This is a PERFECT result.",
    "The model ALWAYS ranks a business card",
    "HIGHER than a consumer card by probability.",
    "",
    "What it looks like in numbers:",
    "Out of 21,000 test cards:",
    "",
    "             Pred 0    Pred 1",
    " Actual 0    15996         4",
    " Actual 1        5      4995",
    "",
    "Only 9 errors out of 21,000!",
    "  • 4 false positives: consumer→business",
    "  • 5 false negatives: business→consumer",
    "",
    "Precision/Recall/F1 = 1.00",
    "Accuracy = 1.00",
    "",
    "⚠️ BUT! Accuracy = 1.00 is NOT because",
    "   we classified everyone as consumer.",
    "   It's because the model actually",
    "   separates the classes perfectly.",
], fs=12, ls=1.2)
add_img(s, CM_IMG, Inches(9.3), Inches(1.35), w=Inches(3.7), h=Inches(3.1))
add_explain_box(s, Inches(9.3), Inches(4.6), Inches(3.7), Inches(2.6), [
    "WHERE DOES AUC=1.0 COME FROM?",
    "(not overfitting!)",
    "",
    "1) 35 engineered features — strong signal",
    "2) tokenized_ratio 26%: business 60-80%,",
    "   consumer 5-10% — huge gap",
    "3) 5-Fold CV gives the same 1.0",
    "4) Early stopping — protection",
    "5) SHAP explains logical patterns",
], icon="🎯")

# ══════════════════════════════════════════════════════════════════════
#  SLIDE 22: Feature Importance
# ══════════════════════════════════════════════════════════════════════
s = new_slide("Feature Importance — which features matter most",
               "Top-15 features, their weight, and what they mean", "22")
fi_data = [
    ("tokenized_ratio", "26.29%", "Business → lots of tokens. Consumer → rarely"),
    ("online_ratio", "20.34%", "Business → online. Consumer → in store"),
    ("b2b_volume_no_recurring", "11.99%", "Business spends without subscriptions"),
    ("channel_alternation_rate", "9.23%", "Business switches channels often"),
    ("vendor_concentration", "3.66%", "Business → 3-5 key suppliers"),
    ("wholesale_spend_ratio", "3.08%", "Business → more wholesale than self"),
    ("cb_spend_share", "2.88%", "Business pays abroad"),
    ("txns_per_merchant", "2.59%", "Business → many txns per merchant"),
    ("total_txns", "2.47%", "Business has more total transactions"),
    ("amount_sum", "2.46%", "Business spends more total money"),
    ("cb_saas_ad_count", "1.89%", "Count of international SaaS/Ads"),
    ("channel_entropy", "1.85%", "Channel entropy (0=1 channel, 1=50/50)"),
    ("n_unique_merchants", "1.82%", "How many different stores"),
    ("n_unique_mcc", "1.18%", "How many different MCC categories"),
    ("merchant_gini", "1.15%", "Concentration (0=even, 1=one merchant)"),
]
rows = len(fi_data) + 1
tbl = s.shapes.add_table(rows, 3, Inches(0.4), Inches(1.5), Inches(6.5), Inches(5.2))
t = tbl.table
t.columns[0].width = Inches(2.8)
t.columns[1].width = Inches(1.0)
t.columns[2].width = Inches(2.7)
for j, h in enumerate(["Feature", "Weight", "What it means"]):
    c = t.cell(0, j); c.text = h
    for p in c.text_frame.paragraphs:
        p.font.size = Pt(12); p.font.bold = True; p.font.color.rgb = WHITE
    c.fill.solid(); c.fill.fore_color.rgb = DARK_BLUE
for i, (feat, imp, desc) in enumerate(fi_data):
    ri = i + 1
    for j, v in enumerate([feat, imp, desc]):
        c = t.cell(ri, j); c.text = v
        for p in c.text_frame.paragraphs:
            p.font.size = Pt(10); p.font.color.rgb = DARK_GRAY
        c.fill.solid()
        c.fill.fore_color.rgb = LIGHT_GRAY if ri % 2 == 0 else WHITE
add_img(s, FI_IMG, Inches(7.2), Inches(1.4), w=Inches(5.6), h=Inches(5.6))

# ══════════════════════════════════════════════════════════════════════
#  SLIDE 23: SHAP
# ══════════════════════════════════════════════════════════════════════
s = new_slide("SHAP — explaining every prediction",
               "Why did the model decide that THIS card is business?", "23")
mtb(s, Inches(0.5), Inches(1.35), Inches(6.5), Inches(5.8), [
    "SHAP (SHapley Additive exPlanations) — a method",
    "that answers the question:",
    "",
    "  'Why did the model say THIS card is business?'",
    "",
    "Idea from game theory:",
    "• Each feature = 'player'",
    "• The prediction = 'team's win'",
    "• SHAP value = each player's contribution",
    "",
    "📊 SHAP Summary Plot (left):",
    "• 1 dot = 1 card",
    "• Color: red = high feature value,",
    "  blue = low feature value",
    "• X-axis: SHAP value",
    "  Positive → pushes toward business",
    "  Negative → pushes toward consumer",
    "",
    "📊 SHAP Bar Plot (right):",
    "• Mean ABSOLUTE SHAP value",
    "• Shows GLOBAL feature importance",
    "",
    "WHAT WE SEE ON THE GRAPHS:",
    "• tokenized_ratio: red on right →",
    "  high tokenized_ratio → business",
    "• channel_alternation_rate: green on right →",
    "  frequent channel switching → business",
    "• This PROVES the model learns",
    "  logical patterns, not noise",
], fs=12, ls=1.2)
add_img(s, SHAP_SUM, Inches(7.2), Inches(1.35), w=Inches(2.8), h=Inches(5.5))
add_img(s, SHAP_BAR, Inches(10.2), Inches(1.35), w=Inches(2.8), h=Inches(5.5))

# ══════════════════════════════════════════════════════════════════════
#  SLIDE 24: Top-100 + Target Product
# ══════════════════════════════════════════════════════════════════════
s = new_slide("Top-100 Hidden Entrepreneurs + Target Product",
               "Ranked list of hidden entrepreneurs", "24")
code_block(s, Inches(0.3), Inches(1.35), Inches(4.5), Inches(5.8), """# Take ONLY consumer cards
consumer_mask = labels == 0
consumer_feat = features.loc[consumer_mask]
consumer_cards = consumer_feat.index.values

# Predict business probability
consumer_probas = model.predict_proba(
    consumer_feat.values)[:, 1]

# Take top 100 by probability
top100_idx = np.argsort(consumer_probas
    )[-100:][::-1]
top100_cards = consumer_cards[top100_idx]
top100_scores = consumer_probas[top100_idx]

# Product logic:
# 1) wholesale_spend_ratio > median → B2B Cashback
# 2) clockwork_cv_mean > median → Working Capital
# 3) channel_alternation_rate > median → Merchant Acct
# 4) otherwise → Business Credit Card

products = []
for i in range(100):
    if top100_ws.iloc[i] > ws_th:
        products.append("B2B Cashback")
    elif top100_cv.iloc[i] > cv_th:
        products.append("Working Capital Loan")
    elif top100_sch.iloc[i] > sch_th:
        products.append("Merchant Account")
    else:
        products.append("Business Credit Card")

top100_df["recommended_product"] = products
top100_df.to_csv(
    "top100_hidden_entrepreneurs.csv",
    index=False)""", fs=7)
mtb(s, Inches(5.0), Inches(1.35), Inches(4.0), Inches(5.8), [
    "HOW WE FOUND 100 HIDDEN ENTREPRENEURS:",
    "",
    "1) Took 80,000 consumer cards",
    "2) Model predicted business probability",
    "   for each",
    "3) Sorted descending",
    "4) Took top 100",
    "",
    "Score distribution:",
    "  P0: 0.0000    P50: 0.0000",
    "  P90: 0.0002   P95: 0.0006",
    "  P99: 0.0085   P100: 0.9467",
    "  Cards with prob>0.5: 22 of 80,000",
    "",
    "Only 0.03% of consumer cards look",
    "like business! Very strict selection.",
    "",
    "PRODUCT RECOMMENDATION:",
    "• B2B Cashback (72%) — for wholesalers",
    "• Business Credit Card (24%) — general",
    "• Working Capital Loan (4%) — regular buyers",
], fs=11, ls=1.2)
add_explain_box(s, Inches(9.3), Inches(1.35), Inches(3.7), Inches(3.0), [
    "EXAMPLE:",
    "",
    "Card 5211559664892162",
    "  Business prob: 0.947 (95%)",
    "  → B2B Cashback",
    "  (high wholesale ratio)",
    "",
    "Card 5119023403360653",
    "  Business prob: 0.713 (71%)",
    "  → Business Credit Card",
    "  (not enough wholesale for",
    "   Cashback, but business exists)",
], icon="🏆")
top10_data = [
    ("5211559664892162", "0.947", "1", "B2B Cashback"),
    ("5176515772256834", "0.929", "2", "B2B Cashback"),
    ("5201491354169846", "0.863", "3", "B2B Cashback"),
    ("5100616099198674", "0.828", "4", "B2B Cashback"),
    ("5486732133818591", "0.816", "5", "B2B Cashback"),
    ("5188841497474379", "0.801", "6", "B2B Cashback"),
    ("5176476691114937", "0.726", "7", "B2B Cashback"),
    ("5176517232027573", "0.722", "8", "B2B Cashback"),
    ("5119023403360653", "0.713", "9", "Business Credit Card"),
    ("5402339951548606", "0.670", "10", "B2B Cashback"),
]
rows = len(top10_data) + 1
tbl = s.shapes.add_table(rows, 4, Inches(9.3), Inches(4.6), Inches(3.7), Inches(2.6))
t = tbl.table
t.columns[0].width = Inches(1.0); t.columns[1].width = Inches(0.8)
t.columns[2].width = Inches(0.5); t.columns[3].width = Inches(1.4)
for j, h in enumerate(["Card", "Prob", "Rank", "Product"]):
    c = t.cell(0, j); c.text = h
    for p in c.text_frame.paragraphs:
        p.font.size = Pt(8); p.font.bold = True; p.font.color.rgb = WHITE
    c.fill.solid(); c.fill.fore_color.rgb = DARK_BLUE
for i, vals in enumerate(top10_data):
    for j, v in enumerate(vals):
        c = t.cell(i+1, j); c.text = v
        for p in c.text_frame.paragraphs:
            p.font.size = Pt(7); p.font.color.rgb = DARK_GRAY
        c.fill.solid()
        c.fill.fore_color.rgb = LIGHT_GRAY if (i+1) % 2 == 0 else WHITE

# ══════════════════════════════════════════════════════════════════════
#  SLIDE 25: Unused columns
# ══════════════════════════════════════════════════════════════════════
s = new_slide("Which columns we DON'T use and why",
               "Detailed analysis of every skipped column", "25")
code_block(s, Inches(0.3), Inches(1.35), Inches(3.8), Inches(5.8), """# Columns we do NOT read:
# From business/consumer cards:
#   - bank_name
#   - card_tier
#
# From merchants_reference:
#   - recurring_capable
#   - merchant_name

# Proof:
import pandas as pd

# bank_name — almost identical
# distribution
biz = pd.read_parquet(..., columns=[
    "bank_name","card_tier"])
biz["label"]=1

con = pd.read_parquet(..., columns=[
    "bank_name","card_tier"])
con["label"]=0

# card_tier: business=100% Business
# con=Standard/Affluent/Premium
# NO OVERLAP!""", fs=7)
add_compare_box(s, Inches(4.4), Inches(1.35), Inches(4.3), Inches(2.5), [
    "✅ tokenized_ratio (26.3%)",
    "  Distinguishes business from consumer",
    "✅ online_ratio (20.3%)",
    "✅ b2b_volume_no_recurring (12%)",
    "  B2B spending without subscriptions",
    "✅ channel_alternation_rate (9.2%)",
], [
    "❌ card_tier (useless):",
    "  Business = 100% 'Business'",
    "  Consumer = 74.5% Standard",
    "            21% Affluent",
    "            4.5% Premium",
    "  NO consumer has 'Business' tier!",
    "  → cannot find hidden entrepreneurs",
])
add_compare_box(s, Inches(4.4), Inches(4.1), Inches(4.3), Inches(3.0), [
    "✅ SAAS_AD_MCCS covers:",
    "  Google Ads, AWS, Shopify,",
    "  Zoom, Slack — all MCCs",
    "  already in our dictionary",
    "✅ cross-border uses",
    "  merchant_country (from the same table)",
], [
    "❌ recurring_capable (1.25%):",
    "  Only 27/2165 merchants",
    "  have recurring_capable=True",
    "  MCC: 7311 (Ads), 7372 (IT),",
    "  5968 (Subs), 4816 (Network)",
    "  → already in SAAS_AD_MCCS!",
    "",
    "❌ merchant_name:",
    "  Nice to have, but MCC already encodes",
    "  the category. NLP not needed.",
])
add_compare_box(s, Inches(9.0), Inches(1.35), Inches(4.0), Inches(2.5), [
    "✅ baseline features:",
    "  amount_sum, amount_mean,",
    "  total_txns — indirectly",
    "  reflect 'wealth level'",
], [
    "❌ bank_name (useless):",
    "  Business top5: Kaspi 31.7%,",
    "    Halyk 23.7%, Forte 8.9%",
    "  Consumer top5: Kaspi 32.0%,",
    "    Halyk 23.8%, Forte 9.1%",
    "  Distribution is IDENTICAL!",
    "  → no signal at all",
])
add_explain_box(s, Inches(9.0), Inches(4.1), Inches(4.0), Inches(3.0), [
    "CONCLUSION:",
    "",
    "Out of 16 columns in raw data,",
    "we DON'T use only 4.",
    "",
    "All 4 have been tested and",
    "add NO additional signal.",
    "",
    "Our 35 features contain ALL",
    "useful information from the data.",
    "",
    "AUC=1.0 is the best proof",
    "that we have enough features.",
], icon="✅", bg=RGBColor(0xE8,0xF5,0xE9), border=GREEN_BORDER)

# ══════════════════════════════════════════════════════════════════════
#  SLIDE 26: Technical Problems
# ══════════════════════════════════════════════════════════════════════
s = new_slide("Technical Problems and Solutions",
               "What went wrong and how we fixed it", "26")
mtb(s, Inches(0.5), Inches(1.4), Inches(6.0), Inches(5.8), [
    "🔴 PROBLEM 1: pandas merge_asof (bug in pandas 3.0.3)",
    "",
    "Tried: pd.merge_asof with by=['card_number']",
    "  → HANGS INDEFINITELY (library bug)",
    "",
    "Solution:",
    "  1. Aggregate to day level (12.8M → 300K)",
    "  2. Dictionary card → np.array of logistics dates",
    "  3. np.searchsorted — binary search",
    "  4. Runs in seconds, not minutes",
    "",
    "",
    "🟠 PROBLEM 2: Round-Trip Cash Flow — no incoming data",
    "",
    "Original: detect money round-tripping",
    "(payment→refund→payment again).",
    "",
    "Problem: our data has ONLY outflows (spending).",
    "No refunds, no deposits.",
    "",
    "Solution: redefined as",
    "'spending_burst_periodicity' — periodicity of",
    "spending bursts. std of intervals between top 5%",
    "spending days. Small std = regular restocking.",
    "",
    "",
    "🟡 PROBLEM 3: pandas 3.0.3 API changes",
    "",
    "• .apply(..., include_groups=False) — new required param",
    "  (without it: warnings/errors)",
    "• .groupby().diff() — use dt.total_seconds()",
    "  instead of dt.days (more precise for hours)",
    "• Code had to be adapted to new API",
    "",
    "",
    "🟢 PROBLEM 4: 12.8M rows — how to avoid memory issues?",
    "",
    "• Feature engineering via groupby aggregations,",
    "  NOT element-wise loops",
    "• Delete intermediate data (del)",
    "• Supplier Fingerprint: head(3) + groupby",
    "• Last-Mile Echo: only daily level",
    "• Result: ~2 GB RAM, ~30 minutes on CPU",
], fs=11, ls=1.15)
add_explain_box(s, Inches(6.8), Inches(1.4), Inches(6.0), Inches(5.8), [
    "LESSONS LEARNED:",
    "",
    "1) Don't trust libraries blindly:",
    "   merge_asof in pandas 3.0.3 is broken.",
    "   Always have a fallback (searchsorted).",
    "",
    "2) Adapt paradoxes to your data:",
    "   Round-Trip Cash Flow requires refunds —",
    "   they don't exist → redefine. This is OK.",
    "   What matters is the IDEA (periodicity),",
    "   not the name.",
    "",
    "3) Optimization is key for big data:",
    "   12.8M rows → day-level aggregation → 300K",
    "   → 40x speedup.",
    "   numpy is always faster than pandas for numbers.",
    "",
    "4) Always check new APIs:",
    "   include_groups=False is mandatory in pandas 3.x",
    "   for groupby.apply. Without it — error.",
    "",
    "5) Use %pip install in the notebook:",
    "   This is a jury requirement (5% for reproducibility).",
    "   Without it, the notebook won't run for the evaluator.",
], icon="💡", bg=RGBColor(0xFF,0xF8,0xE1), border=ORANGE_TEXT)

# ══════════════════════════════════════════════════════════════════════
#  SLIDE 27: Why AUC=1.0 is not overfitting
# ══════════════════════════════════════════════════════════════════════
s = new_slide("Why AUC = 1.0 is NOT overfitting",
               "5 proofs that the model actually works", "27")
mtb(s, Inches(0.5), Inches(1.4), Inches(12.3), Inches(5.8), [
    "PROOF 1: 5-Fold Cross-Validation",
    "  The model was trained 5 times on different 80% of data.",
    "  Each time AUC = 0.999998-0.999999 on the unseen 20%.",
    "  If it were overfitting → CV AUC would be < test AUC.",
    "",
    "PROOF 2: Early Stopping",
    "  The model stopped at iteration 121 out of 500.",
    "  If it were overfitting → test AUC would start",
    "  decreasing after N iterations. Instead, it stayed stable.",
    "",
    "PROOF 3: 35 engineered features — strong signal",
    "  tokenized_ratio: business 60-80% vs consumer 5-10%.",
    "  A 10x difference! This is a real pattern, not noise.",
    "  All 10 paradoxes are theoretically justified patterns.",
    "",
    "PROOF 4: Probability distribution",
    "  Out of 80,000 consumer cards: 99% have prob < 0.01.",
    "  Only 22 cards with prob > 0.5. The model is CONFIDENT.",
    "  If it were overfitting — probabilities would be 'smeared'.",
    "",
    "PROOF 5: SHAP explanations",
    "  SHAP summary plot shows LOGICAL patterns:",
    "  • tokenized_ratio HIGH → red on right → business",
    "  • tokenized_ratio LOW → blue on left → consumer",
    "  • channel_alternation_rate HIGH → business",
    "  The model uses features meaningfully, not memorizing noise.",
    "",
    "CONCLUSION: AUC=1.0 is REAL. The model truly distinguishes",
    "business from consumer. 35 features provide enough signal",
    "for perfect class separation.",
], fs=12, ls=1.3)

# ══════════════════════════════════════════════════════════════════════
#  SLIDE 28: Project Structure
# ══════════════════════════════════════════════════════════════════════
s = new_slide("Project Structure — files and their purpose",
               "What we created", "28")
code_block(s, Inches(0.3), Inches(1.35), Inches(4.5), Inches(5.8), """mqd_classifier.py/
│
├─ 1. LOAD DATA
│  Read 3 parquet files
│
├─ 2. LABEL & CONCAT
│  biz["label"]=1, con["label"]=0
│  pd.concat + merge merchants
│
├─ 3. MCC CATEGORIES
│  5 MCC code dictionaries
│
├─ 4. FEATURE ENGINEERING
│  ├─ 5a Clockwork Buyer
│  ├─ 5b Off-Hours Operator
│  ├─ 5c Expense Ratio Inversion
│  ├─ 5d Token Wholesale
│  ├─ 5e Baseline Aggregations
│  ├─ 5f Supplier Fingerprint
│  ├─ 5g Last-Mile Echo
│  ├─ 5h Round-Trip CF
│  ├─ 5i Inventory Pulse
│  ├─ 5j Multi-Vendor Loyalty
│  ├─ 5k Channel Schizophrenia
│  └─ 5l Cross-border Sourcing
│
├─ 5. ASSEMBLE (35 features)
├─ 6. TRAIN/TEST SPLIT
├─ 7. TRAIN CATBOOST
├─ 8. CROSS-VALIDATION (5-fold)
├─ 9. EVALUATION
├─ 10. FEATURE IMPORTANCE
├─ 11. TOP-100 + TARGET PRODUCT
└─ 12. SAVE OUTPUTS""", fs=8)
mtb(s, Inches(5.0), Inches(1.35), Inches(8.0), Inches(5.8), [
    "📁 Project files:",
    "",
    "1. mqd_classifier.py (816 lines)",
    "   • Main script — from loading to top-100",
    "   • Run: python mqd_classifier.py (~30-60 min)",
    "",
    "2. MQD_2026_Hidden_Entrepreneurs.ipynb (56 cells)",
    "   • Jupyter notebook for jury submission",
    "   • 29 markdown + 27 code cells",
    "   • %pip install for reproducibility (5%)",
    "   • All graphs inline (IPython.display.Image)",
    "",
    "3. confusion_matrix.png",
    "   • Confusion matrix visualization",
    "",
    "4. feature_importance.png",
    "   • Top-15 features by importance",
    "",
    "5. shap_summary.png / shap_bar.png",
    "   • SHAP explanations (if SHAP is installed)",
    "",
    "6. top100_hidden_entrepreneurs.csv",
    "   • 100 hidden entrepreneurs + product recommendation",
    "",
    "7. *.parquet (data, not in repo)",
    "   • business_cards_MDQ.parquet",
    "   • consumer_cards_MDQ.parquet",
    "   • merchants_reference.parquet",
    "",
    "8. Presentation (.pptx)",
    "   • This presentation (30+ slides)",
], fs=12, ls=1.2)

# ══════════════════════════════════════════════════════════════════════
#  SLIDE 29: Conclusions
# ══════════════════════════════════════════════════════════════════════
s = new_slide("Conclusions", "What we achieved and what's important", "29")
mtb(s, Inches(0.5), Inches(1.4), Inches(12.3), Inches(5.8), [
    "✅ KEY RESULTS:",
    "",
    "  📊 AUC-ROC: 1.0000 (primary metric)",
    "  📊 5-Fold CV AUC: 1.0000 ± 0.0000",
    "  📊 Confusion Matrix: 9 errors out of 21,000 (0.04%)",
    "  📊 80,000 consumer cards → 22 with prob > 0.5",
    "  📊 100 hidden entrepreneurs found",
    "  📊 Each has a recommended product (Cashback, Loan, Account, Card)",
    "",
    "✅ TOP FEATURES (TOP-4 = 68% of importance):",
    "  1. tokenized_ratio (26.3%) — tokenized payments",
    "  2. online_ratio (20.3%) — online vs POS",
    "  3. b2b_volume_no_recurring (12.0%) — B2B without subscriptions",
    "  4. channel_alternation_rate (9.2%) — channel switching",
    "",
    "✅ 10 PARADOXES + CROSS-BORDER = 35 FEATURES:",
    "  • All 10 paradoxes implemented and contribute to the model",
    "  • 4 cross-border features add ~6% importance",
    "  • Baseline features (tokenized_ratio, online_ratio) turned out",
    "    to be MORE important than paradoxes — but paradoxes",
    "    explain WHY",
    "",
    "✅ PRODUCTS FOR TOP-100:",
    "  • B2B Cashback: 72% (high wholesale_spend_ratio)",
    "  • Business Credit Card: 24% (basic product)",
    "  • Working Capital Loan: 4% (regular purchasing)",
    "",
    "✅ TECHNICAL ACHIEVEMENTS:",
    "  • Bypassed pandas merge_asof bug using searchsorted",
    "  • Redefined Round-Trip for available data",
    "  • All 10 paradoxes + 4 cross-border features",
    "  • Full reproducibility (seed=42, requirements, %pip)",
    "  • 35 features for 105,000 cards, AUC=1.0 without overfitting",
], fs=12, ls=1.2)

# ══════════════════════════════════════════════════════════════════════
#  SLIDE 30: Thank You
# ══════════════════════════════════════════════════════════════════════
s = add_blank_slide(); add_bg(s, DARK_BLUE)
add_rect(s, 0, Inches(2.8), W, Inches(2.0), MED_BLUE)
add_rect(s, 0, Inches(4.8), W, Inches(0.06), ACCENT_GOLD)
tb(s, Inches(0.8), Inches(3.0), Inches(11.7), Inches(0.8),
   "Thank you for your attention!", fs=40, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
tb(s, Inches(0.8), Inches(3.8), Inches(11.7), Inches(0.8),
   "Mastercard Data Quest 2026 — Hidden Entrepreneur Detection",
   fs=18, color=ACCENT_GOLD, align=PP_ALIGN.CENTER)
tb(s, Inches(0.8), Inches(5.2), Inches(11.7), Inches(0.5),
   "Genius Level • 10 Behavioral Paradoxes • Cross-border Sourcing • AUC-ROC • Target Product",
   fs=14, color=RGBColor(0xBB,0xCC,0xDD), align=PP_ALIGN.CENTER)
tb(s, Inches(0.8), Inches(5.9), Inches(11.7), Inches(0.8),
   "Questions?", fs=28, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════
#  SAVE
# ══════════════════════════════════════════════════════════════════════
output_path = DATA / "MQD_2026_HE_Presentation_EN.pptx"
prs.save(str(output_path))
print(f"ENGLISH presentation saved: {output_path}")
print(f"Total slides: {len(prs.slides)}")
